#!/usr/bin/env python3
"""
PPTX to Video Pipeline
Converts PowerPoint presentations to MP4 video slideshows with voiceover.
"""

import os
import sys
import subprocess
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont
import io


class PPTXToVideoConverter:
    """Main converter class for PPTX to MP4 pipeline."""
    
    def __init__(self, input_dir="input", output_dir="output", temp_dir="temp", background_path=None):
        """
        Initialize the converter.
        
        Args:
            input_dir: Directory containing input files (PPTX and INSTRUKSI.txt)
            output_dir: Directory for output video
            temp_dir: Temporary directory for intermediate files
            background_path: (Deprecated) No longer used - backgrounds are extracted from each slide
        """
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self.temp_dir = Path(temp_dir)
        self.background_path = Path(background_path) if background_path else None
        
        # Create subdirectories
        self.slides_dir = self.temp_dir / "slides"
        self.audio_dir = self.temp_dir / "audio"
        self.videos_dir = self.temp_dir / "videos"
        
        # Ensure directories exist
        for directory in [self.output_dir, self.temp_dir, self.slides_dir, 
                         self.audio_dir, self.videos_dir]:
            directory.mkdir(parents=True, exist_ok=True)
    
    def check_dependencies(self):
        """Check if required external tools are available."""
        # Check for FFmpeg
        try:
            subprocess.run(["ffmpeg", "-version"], 
                         stdout=subprocess.DEVNULL, 
                         stderr=subprocess.DEVNULL, 
                         check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            print("ERROR: FFmpeg is not installed or not in PATH")
            print("Please install FFmpeg: https://ffmpeg.org/download.html")
            sys.exit(1)
        
        # Check for LibreOffice (optional, for better slide rendering)
        try:
            subprocess.run(["soffice", "--version"], 
                         stdout=subprocess.DEVNULL, 
                         stderr=subprocess.DEVNULL, 
                         check=True)
            self.has_libreoffice = True
        except (subprocess.CalledProcessError, FileNotFoundError):
            self.has_libreoffice = False
            print("WARNING: LibreOffice not found. Will use python-pptx for slide rendering.")
    
    def extract_text_from_slide(self, slide):
        """
        Extract all text from a slide.
        
        Args:
            slide: python-pptx slide object
            
        Returns:
            str: All text content from the slide
        """
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        return " ".join(text_parts)
    
    def extract_images_from_slide(self, slide, slide_num):
        """
        Extract images from a slide.
        
        Args:
            slide: python-pptx slide object
            slide_num: Slide number (1-indexed)
            
        Returns:
            list: List of saved image paths
        """
        image_paths = []
        img_count = 0
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    
                    img_count += 1
                    img_filename = f"slide{slide_num:03d}_img{img_count}.{ext}"
                    img_path = self.slides_dir / img_filename
                    
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    
                    image_paths.append(img_path)
                    print(f"  Extracted image: {img_filename}")
                except Exception as e:
                    print(f"  Warning: Could not extract image: {e}")
        
        return image_paths
    
    def extract_background_from_slide(self, slide, prs, width_px, height_px):
        """
        Extract background (image or solid color) from a slide.
        
        Args:
            slide: python-pptx slide object
            prs: python-pptx Presentation object
            width_px: Target width in pixels
            height_px: Target height in pixels
            
        Returns:
            PIL.Image: Background image
        """
        # Default to white background
        bg_color = (255, 255, 255)
        
        try:
            # Try to get background from slide
            bg = slide.background
            fill = bg.fill
            
            # Check if slide follows master or has custom background
            if hasattr(fill, 'type') and fill.type is not None:
                fill_type = fill.type
                
                if fill_type == MSO_FILL_TYPE.SOLID or (hasattr(fill, 'fore_color') and fill.fore_color):
                    try:
                        # Solid color fill
                        if hasattr(fill.fore_color, 'rgb'):
                            rgb = fill.fore_color.rgb
                            bg_color = (rgb[0], rgb[1], rgb[2])
                    except Exception:
                        # If color extraction fails, use default
                        pass
                
                elif fill_type == MSO_FILL_TYPE.PICTURE:
                    # Picture fill backgrounds are not yet supported
                    # Falls through to default white background
                    pass
            
            # If following master background, try to get from master
            if slide.follow_master_background and prs.slide_masters:
                try:
                    master = prs.slide_masters[0]
                    master_bg = master.background
                    master_fill = master_bg.fill
                    
                    if hasattr(master_fill, 'fore_color') and master_fill.fore_color:
                        if hasattr(master_fill.fore_color, 'rgb'):
                            rgb = master_fill.fore_color.rgb
                            bg_color = (rgb[0], rgb[1], rgb[2])
                except Exception:
                    # Master background extraction failed, use default
                    pass
        
        except Exception:
            # If extraction fails, use default white background
            pass
        
        # Create background image with the extracted or default color
        bg_image = Image.new('RGB', (width_px, height_px), bg_color)
        
        return bg_image
    
    def render_text_on_image(self, img, slide, prs):
        """
        Render text shapes from slide onto the image using Pillow.
        
        Args:
            img: PIL Image to draw on
            slide: python-pptx slide object
            prs: python-pptx Presentation object
            
        Returns:
            PIL.Image: Image with text rendered
        """
        draw = ImageDraw.Draw(img)
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        img_width, img_height = img.size
        
        # Scale factors to convert EMU coordinates to pixels
        scale_x = img_width / slide_width
        scale_y = img_height / slide_height
        
        # Try to load a font, fallback to default if not available
        font_large = None
        font_medium = None
        
        try:
            # Try common font paths for different operating systems
            font_configs = [
                # Linux (Debian/Ubuntu)
                {
                    'bold': "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                    'regular': "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
                },
                # Linux (RHEL/Fedora)
                {
                    'bold': "/usr/share/fonts/dejavu-sans-fonts/DejaVuSans-Bold.ttf",
                    'regular': "/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf"
                },
                # Windows
                {
                    'bold': "C:\\Windows\\Fonts\\arialbd.ttf",
                    'regular': "C:\\Windows\\Fonts\\arial.ttf"
                },
                # macOS
                {
                    'bold': "/Library/Fonts/Arial Bold.ttf",
                    'regular': "/Library/Fonts/Arial.ttf"
                },
            ]
            
            for config in font_configs:
                try:
                    font_large = ImageFont.truetype(config['bold'], 48)
                    font_medium = ImageFont.truetype(config['regular'], 32)
                    break
                except (OSError, IOError):
                    continue
            
            # If no font found, raise to use default
            if font_large is None:
                raise Exception("No TrueType font found")
                
        except Exception:
            # Fallback to default font
            font_large = ImageFont.load_default()
            font_medium = ImageFont.load_default()
        
        # Iterate through shapes and draw text
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                try:
                    # Get shape position and size in EMUs
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Convert to pixels
                    x = int(left * scale_x)
                    y = int(top * scale_y)
                    w = int(width * scale_x)
                    h = int(height * scale_y)
                    
                    # Get text
                    text = shape.text.strip()
                    
                    # Choose font based on shape type/position
                    # Title shapes are usually at the top
                    if y < img_height * 0.3:
                        font = font_large
                        color = (0, 0, 0)  # Black
                    else:
                        font = font_medium
                        color = (50, 50, 50)  # Dark gray
                    
                    # Simple text wrapping - split by words
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        test_line = ' '.join(current_line + [word])
                        # Use textbbox instead of deprecated textsize
                        try:
                            bbox = draw.textbbox((0, 0), test_line, font=font)
                            test_width = bbox[2] - bbox[0]
                        except (AttributeError, TypeError):
                            # Fallback for older Pillow versions
                            test_width = len(test_line) * 10
                        
                        if test_width <= w - 20:  # Leave some margin
                            current_line.append(word)
                        else:
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    # Draw each line
                    line_y = y + 10  # Small top margin
                    for line in lines:
                        draw.text((x + 10, line_y), line, font=font, fill=color)
                        # Estimate line height
                        try:
                            bbox = draw.textbbox((0, 0), line, font=font)
                            line_height = bbox[3] - bbox[1]
                        except (AttributeError, TypeError):
                            line_height = 40
                        line_y += line_height + 5  # Line spacing
                        
                        # Stop if we exceed the shape height
                        if line_y > y + h:
                            break
                
                except Exception as e:
                    # If text rendering fails for this shape, continue with others
                    print(f"    Warning: Could not render text for shape: {e}")
                    continue
        
        return img
    
    def convert_slide_to_png_libreoffice(self, pptx_path):
        """
        Convert PPTX slides to PNG using LibreOffice.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            list: List of PNG file paths
        """
        print("Converting slides to PNG using LibreOffice...")
        
        # First convert to PDF
        pdf_path = self.temp_dir / "presentation.pdf"
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(self.temp_dir),
            str(pptx_path)
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True)
        except subprocess.CalledProcessError as e:
            print(f"LibreOffice conversion failed: {e}")
            return None
        
        # Rename the output to our expected name
        generated_pdf = self.temp_dir / f"{pptx_path.stem}.pdf"
        if generated_pdf.exists() and generated_pdf != pdf_path:
            generated_pdf.rename(pdf_path)
        
        # Convert PDF to PNG using pdftoppm or similar
        try:
            cmd = [
                "pdftoppm",
                "-png",
                "-r", "300",
                str(pdf_path),
                str(self.slides_dir / "slide")
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Get list of generated PNGs
            png_files = sorted(self.slides_dir.glob("slide-*.png"))
            
            # Rename to our format (slide001.png, slide002.png, etc.)
            renamed_files = []
            for idx, png_file in enumerate(png_files, 1):
                new_name = self.slides_dir / f"slide{idx:03d}.png"
                png_file.rename(new_name)
                renamed_files.append(new_name)
            
            return renamed_files
        except (subprocess.CalledProcessError, FileNotFoundError):
            print("pdftoppm not available, trying alternative method...")
            return None
    
    def convert_slide_to_png_pythonpptx(self, prs):
        """
        Convert PPTX slides to PNG using python-pptx and Pillow.
        Extracts background from each slide and overlays text.
        
        Args:
            prs: python-pptx Presentation object
            
        Returns:
            list: List of PNG file paths
        """
        print("Converting slides to PNG with dynamic backgrounds using python-pptx...")
        png_files = []
        
        # Get slide dimensions (in EMUs - English Metric Units)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert to pixels (assuming 96 DPI)
        width_px = int(slide_width / 9525)
        height_px = int(slide_height / 9525)
        
        # Standard HD dimensions
        if width_px < 1280 or height_px < 720:
            width_px, height_px = 1920, 1080
        
        for idx, slide in enumerate(prs.slides, 1):
            print(f"  Processing slide {idx}...")
            
            # Extract background from slide
            img = self.extract_background_from_slide(slide, prs, width_px, height_px)
            
            # Render text shapes on top of background
            img = self.render_text_on_image(img, slide, prs)
            
            # Save the rendered slide
            png_path = self.slides_dir / f"slide{idx:03d}.png"
            img.save(png_path)
            png_files.append(png_path)
            print(f"  Created: {png_path.name}")
        
        return png_files
    
    def generate_audio_for_slide(self, text, slide_num, language='en'):
        """
        Generate audio for slide text using gTTS.
        
        Args:
            text: Text to convert to speech
            slide_num: Slide number (1-indexed)
            language: Language code (default: 'en')
            
        Returns:
            Path: Path to generated audio file
        """
        audio_path = self.audio_dir / f"slide{slide_num:03d}.mp3"
        
        if not text or text.strip() == "":
            text = f"Slide {slide_num}"
        
        try:
            tts = gTTS(text=text, lang=language, slow=False)
            tts.save(str(audio_path))
            print(f"  Generated audio: {audio_path.name}")
            return audio_path
        except Exception as e:
            print(f"  Error generating audio: {e}")
            # Create a silent audio file as fallback
            return self.create_silent_audio(audio_path)
    
    def create_silent_audio(self, audio_path, duration=2.0):
        """
        Create a silent audio file using FFmpeg.
        
        Args:
            audio_path: Output path for audio file
            duration: Duration in seconds
            
        Returns:
            Path: Path to generated audio file
        """
        cmd = [
            "ffmpeg",
            "-f", "lavfi",
            "-i", f"anullsrc=r=44100:cl=stereo",
            "-t", str(duration),
            "-q:a", "9",
            "-acodec", "libmp3lame",
            "-y",
            str(audio_path)
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True)
            return audio_path
        except subprocess.CalledProcessError:
            return None
    
    def combine_slide_and_audio(self, slide_path, audio_path, slide_num):
        """
        Combine slide PNG (with background and text already rendered) and audio into a video using FFmpeg.
        
        Args:
            slide_path: Path to slide PNG (already includes background and text)
            audio_path: Path to audio file
            slide_num: Slide number (1-indexed)
            
        Returns:
            Path: Path to generated video file
        """
        video_path = self.videos_dir / f"slide{slide_num:03d}.mp4"
        
        # Slide PNG already contains background and text, just combine with audio
        cmd = [
            "ffmpeg",
            "-loop", "1",
            "-i", str(slide_path),
            "-i", str(audio_path),
            "-c:v", "libx264",
            "-tune", "stillimage",
            "-c:a", "aac",
            "-b:a", "192k",
            "-pix_fmt", "yuv420p",
            "-shortest",
            "-y",
            str(video_path)
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True)
            print(f"  Created video: {video_path.name}")
            return video_path
        except subprocess.CalledProcessError as e:
            print(f"  Error creating video: {e}")
            return None
    
    def concatenate_videos(self, video_paths):
        """
        Concatenate all slide videos into final output.
        
        Args:
            video_paths: List of video file paths
            
        Returns:
            Path: Path to final output video
        """
        output_path = self.output_dir / "output.mp4"
        concat_file = self.temp_dir / "concat_list.txt"
        
        # Create concat file
        with open(concat_file, "w") as f:
            for video_path in video_paths:
                f.write(f"file '{video_path.absolute()}'\n")
        
        cmd = [
            "ffmpeg",
            "-f", "concat",
            "-safe", "0",
            "-i", str(concat_file),
            "-c", "copy",
            "-y",
            str(output_path)
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True)
            print(f"\n✓ Final video created: {output_path}")
            return output_path
        except subprocess.CalledProcessError as e:
            print(f"Error concatenating videos: {e}")
            return None
    
    def process(self, pptx_filename="slides.pptx", language='en'):
        """
        Main processing pipeline.
        
        Args:
            pptx_filename: Name of PPTX file in input directory
            language: Language code for TTS (default: 'en')
        """
        pptx_path = self.input_dir / pptx_filename
        
        if not pptx_path.exists():
            print(f"ERROR: PPTX file not found: {pptx_path}")
            sys.exit(1)
        
        print(f"Processing: {pptx_path}")
        print("=" * 60)
        
        # Note: Background path is deprecated - backgrounds are now extracted from each slide
        if self.background_path:
            print(f"NOTE: --background parameter is deprecated.")
            print("Backgrounds are now automatically extracted from each slide.")
        
        # Check dependencies
        self.check_dependencies()
        
        # Load presentation
        print("\n1. Loading PPTX...")
        prs = Presentation(str(pptx_path))
        num_slides = len(prs.slides)
        print(f"   Found {num_slides} slides")
        
        # Extract text and images from all slides
        print("\n2. Extracting text and images from slides...")
        slide_texts = []
        for idx, slide in enumerate(prs.slides, 1):
            print(f"\n   Slide {idx}:")
            text = self.extract_text_from_slide(slide)
            slide_texts.append(text)
            print(f"   Text: {text[:100]}..." if len(text) > 100 else f"   Text: {text}")
            
            # Extract images (optional)
            images = self.extract_images_from_slide(slide, idx)
        
        # Convert slides to PNG
        print("\n3. Converting slides to PNG images...")
        if self.has_libreoffice:
            png_files = self.convert_slide_to_png_libreoffice(pptx_path)
            if png_files is None:
                png_files = self.convert_slide_to_png_pythonpptx(prs)
        else:
            png_files = self.convert_slide_to_png_pythonpptx(prs)
        
        if not png_files:
            print("ERROR: Failed to convert slides to PNG")
            sys.exit(1)
        
        # Generate audio for each slide
        print("\n4. Generating audio (TTS) for each slide...")
        audio_files = []
        for idx, text in enumerate(slide_texts, 1):
            print(f"\n   Slide {idx}:")
            audio_path = self.generate_audio_for_slide(text, idx, language)
            if audio_path:
                audio_files.append(audio_path)
        
        if len(audio_files) != len(png_files):
            print("ERROR: Mismatch between number of slides and audio files")
            sys.exit(1)
        
        # Combine images and audio into videos
        print("\n5. Combining slides and audio into videos...")
        video_files = []
        for idx, (png_path, audio_path) in enumerate(zip(png_files, audio_files), 1):
            print(f"\n   Slide {idx}:")
            video_path = self.combine_slide_and_audio(png_path, audio_path, idx)
            if video_path:
                video_files.append(video_path)
        
        if not video_files:
            print("ERROR: No videos were created")
            sys.exit(1)
        
        # Concatenate all videos
        print("\n6. Concatenating all slide videos...")
        final_video = self.concatenate_videos(video_files)
        
        if final_video:
            print("\n" + "=" * 60)
            print("✓ PIPELINE COMPLETED SUCCESSFULLY!")
            print(f"✓ Output video: {final_video}")
            print(f"✓ File size: {final_video.stat().st_size / (1024*1024):.2f} MB")
            print("=" * 60)
        else:
            print("\nERROR: Failed to create final video")
            sys.exit(1)


def main():
    """Main entry point."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Convert PPTX presentations to MP4 video slideshows"
    )
    parser.add_argument(
        "--input", "-i",
        default="input",
        help="Input directory containing PPTX file (default: input)"
    )
    parser.add_argument(
        "--output", "-o",
        default="output",
        help="Output directory for video (default: output)"
    )
    parser.add_argument(
        "--temp", "-t",
        default="temp",
        help="Temporary directory for intermediate files (default: temp)"
    )
    parser.add_argument(
        "--pptx", "-p",
        default="slides.pptx",
        help="Name of PPTX file in input directory (default: slides.pptx)"
    )
    parser.add_argument(
        "--language", "-l",
        default="en",
        help="Language code for TTS (default: en, use 'id' for Indonesian)"
    )
    parser.add_argument(
        "--background", "-b",
        default=None,
        help="Path to background PNG image to overlay on slides (default: None)"
    )
    parser.add_argument(
        "--clean",
        action="store_true",
        help="Clean temporary directory before processing"
    )
    
    args = parser.parse_args()
    
    # Clean temp directory if requested
    if args.clean:
        temp_path = Path(args.temp)
        if temp_path.exists():
            print(f"Cleaning temporary directory: {temp_path}")
            shutil.rmtree(temp_path)
    
    # Create converter and process
    converter = PPTXToVideoConverter(
        input_dir=args.input,
        output_dir=args.output,
        temp_dir=args.temp,
        background_path=args.background
    )
    
    converter.process(pptx_filename=args.pptx, language=args.language)


if __name__ == "__main__":
    main()
